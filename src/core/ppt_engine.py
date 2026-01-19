from __future__ import annotations

from pathlib import Path
from typing import Optional
from io import BytesIO
import logging

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


class PptEngineError(RuntimeError):
    """Base error for PPT Engine."""


class SlideNotFoundError(PptEngineError):
    pass


class ShapeNotFoundError(PptEngineError):
    pass


class PptEngine:
    """
    Low-level PowerPoint engine.
    Responsible ONLY for manipulating PPT objects:
    - load/save
    - find slide by slide_key
    - find shape by name
    - set text
    - replace image (BytesIO)

    Does NOT contain any business logic.
    """

    SLIDE_KEY_PREFIX = "SLIDE_KEY_"

    def __init__(self, template_path: Path | str):
        self.template_path = Path(template_path)
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {self.template_path}")

        logger.debug(f"Loading presentation: {self.template_path}")
        self.prs = Presentation(self.template_path)

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    def save(self, output_path: Path | str) -> None:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)

    # ------------------------------------------------------------------
    # Slide handling
    # ------------------------------------------------------------------
    def find_slide_by_key(self, slide_key: str) -> Slide:
        """
        Find slide by invisible anchor shape:
        shape.name == f"SLIDE_KEY_{slide_key}"
        """
        anchor_name = f"{self.SLIDE_KEY_PREFIX}{slide_key}"

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.name == anchor_name:
                    logger.debug(f"Found slide by key '{slide_key}' (slide_id={slide.slide_id})")
                    return slide

        raise SlideNotFoundError(
            f"Slide with key '{slide_key}' not found "
            f"(expected anchor shape name '{anchor_name}')"
        )

    # ------------------------------------------------------------------
    # Shape handling
    # ------------------------------------------------------------------
    def get_shape(self, slide: Slide, shape_name: str) -> BaseShape:
        """
        Get shape by exact name.
        """
        for shape in slide.shapes:
            if shape.name == shape_name:
                return shape

        raise ShapeNotFoundError(
            f"Shape '{shape_name}' not found on slide_id={slide.slide_id}"
        )
        
    # ------------------------------------------------------------------
    # Footer handling
    # ------------------------------------------------------------------
    def set_text_on_layouts(self, shape_name: str, text: str, preserve_format: bool = True) -> int:
        """
        Update a named text shape across ALL slide layouts in ALL slide masters.
        Returns number of updated shapes.
        """
        updated = 0
        for master in self.prs.slide_masters:
            for layout in master.slide_layouts:
                for shape in layout.shapes:
                    if shape.name != shape_name:
                        continue
                    # reuse same implementation style as set_text(), but on a layout shape
                    self._set_text_on_shape(shape, text, preserve_format=preserve_format)
                    updated += 1
        if updated == 0:
            raise ShapeNotFoundError(f"Shape '{shape_name}' not found on any slide layout.")
        return updated
    
    
    def _set_text_on_shape(self, shape, text: str, preserve_format: bool = True) -> None:
        if not shape.has_text_frame:
            raise ValueError(f"Shape '{shape.name}' has no text frame")

        tf = shape.text_frame
        if not preserve_format:
            tf.text = text
            return

        # Ensure at least one paragraph
        if len(tf.paragraphs) == 0:
            tf.text = text
            return

        p0 = tf.paragraphs[0]

        # Ensure at least one run
        if len(p0.runs) == 0:
            p0.text = text
            return

        # Put full text into first run (keep its formatting)
        p0.runs[0].text = text

        # Clear remaining runs in first paragraph (prevents append)
        for r in p0.runs[1:]:
            r.text = ""

        # Clear other paragraphs (if any)
        for p in tf.paragraphs[1:]:
            for r in p.runs:
                r.text = ""


    # ------------------------------------------------------------------
    # Text handling
    # ------------------------------------------------------------------
    def set_text(self, slide: Slide, shape_name: str, text: str, preserve_format: bool = True) -> None:
        """
        Set text on a text-capable shape.

        If preserve_format=True (default):
        - Keep existing paragraph/run formatting by replacing text in the first run,
            and clearing other runs' text.

        If preserve_format=False:
        - Clear the text frame and set plain text (may lose formatting).
        """
        shape = self.get_shape(slide, shape_name)

        if not shape.has_text_frame:
            raise PptEngineError(f"Shape '{shape_name}' has no text frame")

        tf = shape.text_frame

        if not preserve_format:
            tf.clear()
            tf.paragraphs[0].text = text
            logger.debug(f"Set text (no format) on shape '{shape_name}': {text}")
            return

        # --- Preserve formatting path ---
        # Ensure there's at least one paragraph
        if len(tf.paragraphs) == 0:
            tf.text = text
            logger.debug(f"Set text (fallback preserve) on shape '{shape_name}': {text}")
            return

        p0 = tf.paragraphs[0]

        # Ensure there's at least one run in the first paragraph
        if len(p0.runs) == 0:
            # Adding text via paragraph.text may create a run with default formatting,
            # but it's the best we can do if template has no runs.
            p0.text = text
            logger.debug(f"Set text (created run) on shape '{shape_name}': {text}")
            return

        # Put all text into the first run to keep its formatting
        p0.runs[0].text = text

        # Clear remaining runs in the first paragraph
        for r in p0.runs[1:]:
            r.text = ""

        # Clear other paragraphs' runs (keep structure but empty)
        for p in tf.paragraphs[1:]:
            for r in p.runs:
                r.text = ""

        logger.debug(f"Set text (preserve format) on shape '{shape_name}': {text}")


    # ------------------------------------------------------------------
    # Image handling
    # ------------------------------------------------------------------
    def replace_image(
        self,
        slide: Slide,
        shape_name: str,
        image_stream: BytesIO,
    ) -> None:
        """
        Replace an existing picture shape with a new image,
        preserving position and size.
        """
        shape = self.get_shape(slide, shape_name)

        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise PptEngineError(
                f"Shape '{shape_name}' is not a picture (type={shape.shape_type})"
            )

        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        z_order = shape._element.getparent().index(shape._element)

        # Remove old shape
        sp_tree = slide.shapes._spTree
        sp_tree.remove(shape._element)

        # Add new image
        pic = slide.shapes.add_picture(
            image_stream,
            left=left,
            top=top,
            width=width,
            height=height,
        )

        # Try to restore z-order (best effort)
        try:
            sp_tree.remove(pic._element)
            sp_tree.insert(z_order, pic._element)
        except Exception:
            # z-order restore is not critical
            pass

        logger.debug(f"Replaced image on shape '{shape_name}'")

